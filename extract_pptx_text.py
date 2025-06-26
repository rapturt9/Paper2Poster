import json
from pptx import Presentation

prs = Presentation('poster_output.pptx')

for slide_num, slide in enumerate(prs.slides):
    print(f"Slide {slide_num + 1}:")
    for shape_num, shape in enumerate(slide.shapes):
        if hasattr(shape, 'text_frame'):
            text = shape.text_frame.text
            print(f"  Shape {shape_num + 1}: {text[:100]}{'...' if len(text) > 100 else ''}")
    print()