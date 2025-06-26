import os
import sys
import json
import requests
import io
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import pdf2image
from pdfminer.high_level import extract_text

# Load environment variables
load_dotenv()

# Get the OpenRouter API key
openrouter_api_key = os.getenv("OPENROUTER_API_KEY")
if not openrouter_api_key:
    print("Error: OPENROUTER_API_KEY environment variable is not set.")
    sys.exit(1)

def extract_text_from_pdf(pdf_path):
    """Extract text from the PDF using pdfminer."""
    try:
        # Extract text from PDF
        text = extract_text(pdf_path)
        
        # Limit text length to avoid token limits
        max_length = 10000
        if len(text) > max_length:
            text = text[:max_length] + "... [text truncated due to length]"
        
        print(f"Extracted {len(text)} characters from PDF")
        return text
    except Exception as e:
        print(f"Error extracting text from PDF: {e}")
        return "Error extracting text from PDF"

def call_openrouter_api(prompt, pdf_text):
    """Call the OpenRouter API with GPT-4o-mini to generate poster content."""
    headers = {
        "Authorization": f"Bearer {openrouter_api_key}",
        "Content-Type": "application/json",
        "HTTP-Referer": "https://all-hands.dev",
        "X-Title": "Paper2Poster"
    }
    
    data = {
        "model": "openai/gpt-4o-mini",
        "messages": [
            {
                "role": "system", 
                "content": "You are an expert at creating academic posters from research papers. Extract the key information and organize it into a clear, visually appealing academic poster structure."
            },
            {
                "role": "user", 
                "content": f"Create an academic poster from this research paper. Extract the title, authors, abstract, introduction, methodology, results, and conclusion. Format it as a JSON with the following structure: {{\"title\": \"\", \"authors\": \"\", \"abstract\": \"\", \"introduction\": \"\", \"methodology\": \"\", \"results\": \"\", \"conclusion\": \"\", \"references\": \"\"}}. Here's the paper content:\n\n{pdf_text}\n\n{prompt}"
            }
        ],
        "max_tokens": 2000,
        "response_format": {"type": "json_object"}
    }
    
    response = requests.post("https://openrouter.ai/api/v1/chat/completions", headers=headers, json=data)
    
    if response.status_code == 200:
        content = response.json()["choices"][0]["message"]["content"]
        print(f"API Response: {content[:100]}...")  # Print the first 100 chars for debugging
        return content
    else:
        print(f"Error calling OpenRouter API: {response.status_code}")
        print(response.text)
        return None

def create_poster(content, output_path):
    """Create a PowerPoint poster from the generated content."""
    try:
        # Parse the JSON content
        poster_data = json.loads(content)
        
        # Create a new presentation
        prs = Presentation()
        
        # Set slide dimensions for a poster (48x36 inches)
        prs.slide_width = Inches(48)
        prs.slide_height = Inches(36)
        
        # Add a slide
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(2), Inches(1), Inches(44), Inches(3))
        title_frame = title_box.text_frame
        title_para = title_frame.add_paragraph()
        title_para.text = poster_data["title"]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(72)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(0, 0, 128)  # Dark blue
        
        # Add authors
        authors_box = slide.shapes.add_textbox(Inches(2), Inches(4), Inches(44), Inches(1))
        authors_frame = authors_box.text_frame
        authors_para = authors_frame.add_paragraph()
        authors_para.text = poster_data["authors"]
        authors_para.alignment = PP_ALIGN.CENTER
        authors_para.font.size = Pt(40)
        authors_para.font.italic = True
        
        # Add abstract
        abstract_box = slide.shapes.add_textbox(Inches(2), Inches(6), Inches(44), Inches(4))
        abstract_frame = abstract_box.text_frame
        abstract_title = abstract_frame.add_paragraph()
        abstract_title.text = "Abstract"
        abstract_title.font.size = Pt(48)
        abstract_title.font.bold = True
        abstract_content = abstract_frame.add_paragraph()
        abstract_content.text = poster_data["abstract"]
        abstract_content.font.size = Pt(32)
        
        # Create a 2x2 grid for the remaining sections
        # Introduction (top left)
        intro_box = slide.shapes.add_textbox(Inches(2), Inches(11), Inches(21), Inches(10))
        intro_frame = intro_box.text_frame
        intro_title = intro_frame.add_paragraph()
        intro_title.text = "Introduction"
        intro_title.font.size = Pt(48)
        intro_title.font.bold = True
        intro_content = intro_frame.add_paragraph()
        intro_content.text = poster_data["introduction"]
        intro_content.font.size = Pt(28)
        
        # Methodology (top right)
        method_box = slide.shapes.add_textbox(Inches(25), Inches(11), Inches(21), Inches(10))
        method_frame = method_box.text_frame
        method_title = method_frame.add_paragraph()
        method_title.text = "Methodology"
        method_title.font.size = Pt(48)
        method_title.font.bold = True
        method_content = method_frame.add_paragraph()
        method_content.text = poster_data["methodology"]
        method_content.font.size = Pt(28)
        
        # Results (bottom left)
        results_box = slide.shapes.add_textbox(Inches(2), Inches(22), Inches(21), Inches(10))
        results_frame = results_box.text_frame
        results_title = results_frame.add_paragraph()
        results_title.text = "Results"
        results_title.font.size = Pt(48)
        results_title.font.bold = True
        results_content = results_frame.add_paragraph()
        results_content.text = poster_data["results"]
        results_content.font.size = Pt(28)
        
        # Conclusion (bottom right)
        concl_box = slide.shapes.add_textbox(Inches(25), Inches(22), Inches(21), Inches(10))
        concl_frame = concl_box.text_frame
        concl_title = concl_frame.add_paragraph()
        concl_title.text = "Conclusion"
        concl_title.font.size = Pt(48)
        concl_title.font.bold = True
        concl_content = concl_frame.add_paragraph()
        concl_content.text = poster_data["conclusion"]
        concl_content.font.size = Pt(28)
        
        # References at the bottom
        ref_box = slide.shapes.add_textbox(Inches(2), Inches(33), Inches(44), Inches(2))
        ref_frame = ref_box.text_frame
        ref_title = ref_frame.add_paragraph()
        ref_title.text = "References"
        ref_title.font.size = Pt(36)
        ref_title.font.bold = True
        ref_content = ref_frame.add_paragraph()
        ref_content.text = poster_data["references"]
        ref_content.font.size = Pt(20)
        
        # Save the presentation
        prs.save(output_path)
        print(f"Poster saved to {output_path}")
        return True
    except Exception as e:
        print(f"Error creating poster: {e}")
        return False

def main():
    if len(sys.argv) < 2:
        print("Usage: python simple_poster.py <path_to_paper.pdf>")
        sys.exit(1)
    
    pdf_path = sys.argv[1]
    output_path = "poster_output.pptx"
    
    # Extract text from PDF
    print("Extracting text from PDF...")
    pdf_text = extract_text_from_pdf(pdf_path)
    
    # Call OpenRouter API
    print("Generating poster content with GPT-4o-mini...")
    prompt = "Create an academic poster from this research paper. Focus on the main contributions and findings."
    content = call_openrouter_api(prompt, pdf_text)
    
    if content:
        # Create poster
        print("Creating poster...")
        success = create_poster(content, output_path)
        
        if success:
            print(f"Poster successfully created at {output_path}")
        else:
            print("Failed to create poster")
    else:
        print("Failed to generate poster content")

if __name__ == "__main__":
    main()