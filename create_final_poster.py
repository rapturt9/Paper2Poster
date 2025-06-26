import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Load the extracted content
with open('extracted_poster_content.json', 'r') as f:
    content = json.load(f)

# Create a presentation
prs = Presentation()

# Set slide dimensions to a standard poster size (36x24 inches)
prs.slide_width = Inches(36)
prs.slide_height = Inches(24)

# Add a slide
slide_layout = prs.slide_layouts[6]  # blank layout
slide = prs.slides.add_slide(slide_layout)

# Define colors
PRIMARY_COLOR = RGBColor(25, 25, 112)  # Midnight Blue
SECONDARY_COLOR = RGBColor(70, 130, 180)  # Steel Blue
ACCENT_COLOR = RGBColor(0, 0, 128)  # Navy
TEXT_COLOR = RGBColor(0, 0, 0)  # Black
BACKGROUND_COLOR = RGBColor(240, 248, 255)  # Alice Blue

# Set background color
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = BACKGROUND_COLOR

# Add a header bar
header = slide.shapes.add_shape(
    1,  # Rectangle
    Inches(0), Inches(0), 
    Inches(36), Inches(3)
)
header_fill = header.fill
header_fill.solid()
header_fill.fore_color.rgb = PRIMARY_COLOR

# Add title
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(34), Inches(2))
title_frame = title_box.text_frame
title_frame.word_wrap = True
title_para = title_frame.paragraphs[0]
title_para.alignment = PP_ALIGN.CENTER
title_run = title_para.add_run()
title_run.text = content["title"].strip()
title_run.font.size = Pt(60)
title_run.font.bold = True
title_run.font.color.rgb = RGBColor(255, 255, 255)  # White

# Add authors
authors_box = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(34), Inches(1))
authors_frame = authors_box.text_frame
authors_para = authors_frame.paragraphs[0]
authors_para.alignment = PP_ALIGN.CENTER
authors_run = authors_para.add_run()
authors_run.text = content["authors"].strip()
authors_run.font.size = Pt(36)
authors_run.font.italic = True
authors_run.font.color.rgb = TEXT_COLOR

# Function to add a section with a colored header
def add_section(title, content_text, left, top, width, height, color=SECONDARY_COLOR):
    # Add section box with border and rounded corners
    section_box = slide.shapes.add_shape(
        1,  # Rectangle
        left, top, 
        width, height
    )
    section_box.fill.solid()
    section_box.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White
    section_box.line.color.rgb = color
    section_box.line.width = Pt(3)
    
    # Add section header
    header_box = slide.shapes.add_shape(
        1,  # Rectangle
        left, top, 
        width, Inches(1)
    )
    header_box.fill.solid()
    header_box.fill.fore_color.rgb = color
    header_box.line.color.rgb = color
    
    # Add title
    title_box = slide.shapes.add_textbox(
        left + Inches(0.2), 
        top + Inches(0.2), 
        width - Inches(0.4), 
        Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_run = title_para.add_run()
    title_run.text = title
    title_run.font.size = Pt(40)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(255, 255, 255)  # White
    
    # Add content
    content_box = slide.shapes.add_textbox(
        left + Inches(0.5), 
        top + Inches(1.2), 
        width - Inches(1), 
        height - Inches(1.5)
    )
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    content_para = content_frame.paragraphs[0]
    content_run = content_para.add_run()
    content_run.text = content_text.replace(title, "").strip()
    content_run.font.size = Pt(28)
    content_run.font.color.rgb = TEXT_COLOR

# Layout sections
# Abstract (top left)
abstract_box = add_section(
    "Abstract", 
    content["abstract"], 
    Inches(1), Inches(4.5), 
    Inches(10), Inches(6),
    SECONDARY_COLOR
)

# Introduction (top center)
intro_box = add_section(
    "Introduction", 
    content["introduction"], 
    Inches(12), Inches(4.5), 
    Inches(12), Inches(6),
    SECONDARY_COLOR
)

# Methodology (top right)
method_box = add_section(
    "Methodology", 
    content["methodology"], 
    Inches(25), Inches(4.5), 
    Inches(10), Inches(6),
    SECONDARY_COLOR
)

# Results (bottom left)
results_box = add_section(
    "Results", 
    content["results"], 
    Inches(1), Inches(11.5), 
    Inches(16), Inches(6),
    ACCENT_COLOR
)

# Conclusion (bottom center)
conclusion_box = add_section(
    "Conclusion", 
    content["conclusion"], 
    Inches(18), Inches(11.5), 
    Inches(17), Inches(6),
    ACCENT_COLOR
)

# Add a footer bar
footer = slide.shapes.add_shape(
    1,  # Rectangle
    Inches(0), Inches(18.5), 
    Inches(36), Inches(5.5)
)
footer_fill = footer.fill
footer_fill.solid()
footer_fill.fore_color.rgb = PRIMARY_COLOR

# References (bottom)
references_title_box = slide.shapes.add_textbox(
    Inches(1), Inches(19), 
    Inches(34), Inches(1)
)
references_title_frame = references_title_box.text_frame
references_title_para = references_title_frame.paragraphs[0]
references_title_para.alignment = PP_ALIGN.CENTER
references_title_run = references_title_para.add_run()
references_title_run.text = "References"
references_title_run.font.size = Pt(40)
references_title_run.font.bold = True
references_title_run.font.color.rgb = RGBColor(255, 255, 255)  # White

references_box = slide.shapes.add_textbox(
    Inches(1), Inches(20), 
    Inches(34), Inches(3)
)
references_frame = references_box.text_frame
references_para = references_frame.paragraphs[0]
references_para.alignment = PP_ALIGN.CENTER
references_run = references_para.add_run()
references_run.text = content["references"].replace("References", "").strip()
references_run.font.size = Pt(28)
references_run.font.color.rgb = RGBColor(255, 255, 255)  # White

# Save the presentation
output_path = "final_poster.pptx"
prs.save(output_path)
print(f"Final poster saved to {output_path}")