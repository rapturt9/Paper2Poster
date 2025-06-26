import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Load the extracted content
with open('extracted_poster_content.json', 'r') as f:
    content = json.load(f)

# Create a presentation
prs = Presentation()

# Set slide dimensions to a standard poster size (36x24 inches)
prs.slide_width = Inches(36)
prs.slide_height = Inches(24)

# Add a slide
slide_layout = prs.slide_layouts[6]  # blank layout
slide = prs.slides.add_slide(slide_layout)

# Define colors
TITLE_COLOR = RGBColor(0, 0, 128)  # Navy Blue
SECTION_COLOR = RGBColor(0, 51, 102)  # Dark Blue
TEXT_COLOR = RGBColor(0, 0, 0)  # Black
BACKGROUND_COLOR = RGBColor(255, 255, 255)  # White

# Set background color
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = BACKGROUND_COLOR

# Add title
title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(34), Inches(2))
title_frame = title_box.text_frame
title_frame.word_wrap = True
title_para = title_frame.paragraphs[0]
title_para.alignment = PP_ALIGN.CENTER
title_run = title_para.add_run()
title_run.text = content["title"].strip()
title_run.font.size = Pt(60)
title_run.font.bold = True
title_run.font.color.rgb = TITLE_COLOR

# Add authors
authors_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(34), Inches(1))
authors_frame = authors_box.text_frame
authors_para = authors_frame.paragraphs[0]
authors_para.alignment = PP_ALIGN.CENTER
authors_run = authors_para.add_run()
authors_run.text = content["authors"].strip()
authors_run.font.size = Pt(36)
authors_run.font.italic = True
authors_run.font.color.rgb = TEXT_COLOR

# Function to add a section
def add_section(title, content_text, left, top, width, height):
    # Add section box with border
    section_box = slide.shapes.add_textbox(left, top, width, height)
    
    # Add title
    text_frame = section_box.text_frame
    text_frame.word_wrap = True
    title_para = text_frame.paragraphs[0]
    title_run = title_para.add_run()
    title_run.text = title
    title_run.font.size = Pt(40)
    title_run.font.bold = True
    title_run.font.color.rgb = SECTION_COLOR
    
    # Add content
    content_para = text_frame.add_paragraph()
    content_para.space_before = Pt(12)
    content_run = content_para.add_run()
    content_run.text = content_text.replace(title, "").strip()
    content_run.font.size = Pt(28)
    content_run.font.color.rgb = TEXT_COLOR
    
    return section_box

# Layout sections
# Abstract (top left)
abstract_box = add_section(
    "Abstract", 
    content["abstract"], 
    Inches(1), Inches(4.5), 
    Inches(10), Inches(6)
)

# Introduction (top center)
intro_box = add_section(
    "Introduction", 
    content["introduction"], 
    Inches(12), Inches(4.5), 
    Inches(12), Inches(6)
)

# Methodology (top right)
method_box = add_section(
    "Methodology", 
    content["methodology"], 
    Inches(25), Inches(4.5), 
    Inches(10), Inches(6)
)

# Results (bottom left)
results_box = add_section(
    "Results", 
    content["results"], 
    Inches(1), Inches(11.5), 
    Inches(16), Inches(6)
)

# Conclusion (bottom center)
conclusion_box = add_section(
    "Conclusion", 
    content["conclusion"], 
    Inches(18), Inches(11.5), 
    Inches(17), Inches(6)
)

# References (bottom)
references_box = add_section(
    "References", 
    content["references"], 
    Inches(1), Inches(18.5), 
    Inches(34), Inches(4)
)

# Save the presentation
output_path = "enhanced_poster.pptx"
prs.save(output_path)
print(f"Enhanced poster saved to {output_path}")