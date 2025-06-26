import json
from pptx import Presentation

prs = Presentation('poster_output.pptx')
poster_content = {
    "title": "",
    "authors": "",
    "abstract": "",
    "introduction": "",
    "methodology": "",
    "results": "",
    "conclusion": "",
    "references": ""
}

for slide in prs.slides:
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame'):
            text = shape.text_frame.text
            
            # Extract content based on section headers
            if "Evaluating LLM Agent" in text:
                poster_content["title"] = text
            elif "Anonymous Authors" in text:
                poster_content["authors"] = text
            elif text.startswith("Abstract"):
                poster_content["abstract"] = text
            elif text.startswith("Introduction"):
                poster_content["introduction"] = text
            elif text.startswith("Methodology"):
                poster_content["methodology"] = text
            elif text.startswith("Results"):
                poster_content["results"] = text
            elif text.startswith("Conclusion"):
                poster_content["conclusion"] = text
            elif text.startswith("References"):
                poster_content["references"] = text
            
            # Also check for paragraphs
            for paragraph in shape.text_frame.paragraphs:
                paragraph_text = paragraph.text
                if paragraph_text.startswith("Abstract") or "Abstract" in paragraph_text and len(poster_content["abstract"]) == 0:
                    poster_content["abstract"] = shape.text_frame.text
                elif paragraph_text.startswith("Introduction") or "Introduction" in paragraph_text and len(poster_content["introduction"]) == 0:
                    poster_content["introduction"] = shape.text_frame.text
                elif paragraph_text.startswith("Methodology") or "Methodology" in paragraph_text and len(poster_content["methodology"]) == 0:
                    poster_content["methodology"] = shape.text_frame.text
                elif paragraph_text.startswith("Results") or "Results" in paragraph_text and len(poster_content["results"]) == 0:
                    poster_content["results"] = shape.text_frame.text
                elif paragraph_text.startswith("Conclusion") or "Conclusion" in paragraph_text and len(poster_content["conclusion"]) == 0:
                    poster_content["conclusion"] = shape.text_frame.text
                elif paragraph_text.startswith("References") or "References" in paragraph_text and len(poster_content["references"]) == 0:
                    poster_content["references"] = shape.text_frame.text

# Save the extracted content to a JSON file
with open('extracted_poster_content.json', 'w') as f:
    json.dump(poster_content, f, indent=2)

# Print the extracted content
for section, content in poster_content.items():
    print(f"{section.upper()}:")
    print(content[:200] + "..." if len(content) > 200 else content)
    print()